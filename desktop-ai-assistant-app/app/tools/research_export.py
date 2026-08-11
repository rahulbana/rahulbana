"""Save research/results to a file in many formats.

Supported: txt, md, json, xml, csv, xlsx (excel), docx, pptx (ppt), pdf.
"""
from __future__ import annotations

import csv as csv_module
import json as json_module
import re
from datetime import datetime
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape

from ..config import config
from .base import Tool, err, ok


def _safe_name(title: str) -> str:
    slug = re.sub(r"[^\w\s-]", "", title).strip().replace(" ", "_")
    slug = slug[:60] or "research"
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{slug}_{stamp}"


def _load_json(content: str):
    """Return the parsed JSON value if content is valid JSON, else None."""
    if not isinstance(content, str):
        return content
    stripped = content.strip()
    if not stripped or stripped[0] not in "[{":
        return None
    try:
        return json_module.loads(stripped)
    except (json_module.JSONDecodeError, ValueError):
        return None


def _as_records(content: str):
    """If content is a JSON array of objects (or a single object), return a list
    of dicts suitable for a table; otherwise None."""
    data = _load_json(content)
    if isinstance(data, list) and data and all(isinstance(x, dict) for x in data):
        return data
    if isinstance(data, dict):
        return [data]
    return None


def _headers(records: list[dict]) -> list[str]:
    """Union of keys across records, preserving first-seen order."""
    seen: dict[str, None] = {}
    for rec in records:
        for k in rec:
            seen.setdefault(k, None)
    return list(seen)


def _cell(value) -> str:
    """Render a value as a single cell (lists/dicts flattened readably)."""
    if isinstance(value, list):
        return "; ".join(_cell(v) for v in value)
    if isinstance(value, dict):
        return json_module.dumps(value, ensure_ascii=False)
    if value is None:
        return ""
    return str(value)


# --------------------------- format writers ---------------------------

def _save_txt(title: str, content: str, path: Path) -> None:
    path.write_text(f"{title}\n{'=' * len(title)}\n\n{content}\n", encoding="utf-8")


def _save_md(title: str, content: str, path: Path) -> None:
    path.write_text(f"# {title}\n\n{content}\n", encoding="utf-8")


def _save_json(title: str, content: str, path: Path) -> None:
    # If the content is itself JSON, embed the parsed structure so we don't
    # double-encode it into a single string full of escaped newlines.
    parsed = _load_json(content)
    payload = {
        "title": title,
        "content": parsed if parsed is not None else content,
        "generated_at": datetime.now().isoformat(timespec="seconds"),
    }
    path.write_text(json_module.dumps(payload, indent=2, ensure_ascii=False),
                    encoding="utf-8")


def _save_xml(title: str, content: str, path: Path) -> None:
    body = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        "<document>\n"
        f"  <title>{xml_escape(title)}</title>\n"
        f"  <generated_at>{datetime.now().isoformat(timespec='seconds')}</generated_at>\n"
        "  <content>\n"
    )
    for para in content.split("\n\n"):
        para = para.strip()
        if para:
            body += f"    <paragraph>{xml_escape(para)}</paragraph>\n"
    body += "  </content>\n</document>\n"
    path.write_text(body, encoding="utf-8")


def _save_csv(title: str, content: str, path: Path) -> None:
    records = _as_records(content)
    with path.open("w", newline="", encoding="utf-8") as f:
        writer = csv_module.writer(f)
        if records:
            # Tabular JSON content -> a real CSV table.
            headers = _headers(records)
            writer.writerow(headers)
            for rec in records:
                writer.writerow([_cell(rec.get(h)) for h in headers])
        else:
            writer.writerow(["Title", title])
            writer.writerow([])
            writer.writerow(["Line", "Text"])
            for i, line in enumerate(content.splitlines(), start=1):
                writer.writerow([i, line])


def _save_xlsx(title: str, content: str, path: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    records = _as_records(content)

    if records:
        # Tabular JSON content -> a real Excel table with a bold header row.
        headers = _headers(records)
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)
        for rec in records:
            ws.append([_cell(rec.get(h)) for h in headers])
        for i, h in enumerate(headers, start=1):
            width = max(len(str(h)), *(len(_cell(r.get(h))) for r in records)) + 2
            ws.column_dimensions[ws.cell(row=1, column=i).column_letter].width = \
                min(width, 60)
    else:
        ws["A1"] = title
        ws["A1"].font = Font(bold=True, size=14)
        row = 3
        for line in content.splitlines():
            ws.cell(row=row, column=1, value=line)
            row += 1
        ws.column_dimensions["A"].width = 100
    wb.save(path)


def _save_docx(title: str, content: str, path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for block in content.split("\n\n"):
        block = block.strip()
        if not block:
            continue
        if block.startswith("# "):
            doc.add_heading(block[2:], level=1)
        elif block.startswith("## "):
            doc.add_heading(block[3:], level=2)
        else:
            doc.add_paragraph(block)
    doc.save(path)


def _save_pptx(title: str, content: str, path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    # Title slide.
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = datetime.now().strftime("%B %d, %Y")

    # One content slide per paragraph block (keeps slides readable).
    blocks = [b.strip() for b in content.split("\n\n") if b.strip()] or [content]
    for block in blocks:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        heading = block.splitlines()[0][:80]
        slide.shapes.title.text = heading
        body = slide.placeholders[1].text_frame
        body.text = block
        for para in body.paragraphs:
            for run in para.runs:
                run.font.size = Pt(16)
    prs.save(path)


def _save_pdf(title: str, content: str, path: Path) -> None:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    def latin1(s: str) -> str:
        return s.encode("latin-1", "replace").decode("latin-1")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, latin1(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(2)
    pdf.set_font("Helvetica", size=12)
    for line in content.split("\n"):
        text = latin1(line)
        if not text.strip():
            pdf.ln(4)
            continue
        pdf.multi_cell(0, 7, text, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.output(str(path))


# format alias -> (writer, extension)
_WRITERS = {
    "txt": (_save_txt, ".txt"),
    "text": (_save_txt, ".txt"),
    "md": (_save_md, ".md"),
    "markdown": (_save_md, ".md"),
    "json": (_save_json, ".json"),
    "xml": (_save_xml, ".xml"),
    "csv": (_save_csv, ".csv"),
    "xlsx": (_save_xlsx, ".xlsx"),
    "excel": (_save_xlsx, ".xlsx"),
    "docx": (_save_docx, ".docx"),
    "doc": (_save_docx, ".docx"),
    "pptx": (_save_pptx, ".pptx"),
    "ppt": (_save_pptx, ".pptx"),
    "pdf": (_save_pdf, ".pdf"),
}

# Canonical formats surfaced to the model/UI.
FORMATS = ["txt", "md", "json", "xml", "csv", "xlsx", "docx", "pptx", "pdf"]


def save_research(title: str, content: str, format: str = "md") -> dict:
    """Save research content to a file in the requested format."""
    fmt = format.lower().strip()
    if fmt not in _WRITERS:
        return err(f"Unsupported format '{format}'. Use one of: {', '.join(FORMATS)}.")

    writer, ext = _WRITERS[fmt]
    out = config.EXPORT_DIR / (_safe_name(title) + ext)
    try:
        writer(title, content, out)
    except Exception as exc:
        return err(f"Failed to save {fmt} file: {exc}")

    return ok({
        "title": title,
        "format": fmt,
        "path": str(out),
        "bytes": out.stat().st_size,
    }, message=f"Saved research to {out}")


def get_tools() -> list[Tool]:
    return [
        Tool(
            name="save_research",
            description="Save research or results to a file on disk. Supported "
                        "formats: txt, md, json, xml, csv, xlsx (Excel), docx "
                        "(Word), pptx (PowerPoint), pdf. Returns the saved path.",
            category="Research",
            parameters={
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Document title."},
                    "content": {"type": "string", "description": "The body content to save."},
                    "format": {
                        "type": "string",
                        "enum": FORMATS,
                        "description": "Output format. Default 'md'.",
                    },
                },
                "required": ["title", "content"],
            },
            func=save_research,
        )
    ]
