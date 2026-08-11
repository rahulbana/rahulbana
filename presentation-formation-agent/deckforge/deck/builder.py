"""Render a :class:`~deckforge.models.Deck` into a polished .pptx file.

The builder is pure presentation logic — it takes validated data and a
:class:`~deckforge.deck.themes.Theme` and draws 16:9 slides with accent
shapes, consistent typography and per-layout compositions. No LLM or network
access happens here, which keeps rendering fast and deterministic.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..models import Deck, Slide, SlideLayout
from .themes import Theme, get_theme

# 16:9 canvas
_W = Inches(13.333)
_H = Inches(7.5)
_MARGIN = Inches(0.9)


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 == blank


def _fill_background(slide, color: RGBColor) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _rect(slide, x, y, w, h, color: RGBColor, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(
    slide,
    x,
    y,
    w,
    h,
    text: str,
    *,
    font: str,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.1,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return box


def _bullets(slide, x, y, w, h, items, theme: Theme, *, size: int = 18):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        # Accent bullet glyph
        marker = p.add_run()
        marker.text = "▸  "  # ▸
        marker.font.name = theme.body_font
        marker.font.size = Pt(size)
        marker.font.color.rgb = theme.primary
        marker.font.bold = True
        body = p.add_run()
        body.text = item
        body.font.name = theme.body_font
        body.font.size = Pt(size)
        body.font.color.rgb = theme.text
    return box


def _accent_corner(slide, theme: Theme) -> None:
    """A decorative diagonal accent in the top-left + a thin base bar."""
    _rect(slide, Emu(0), Emu(0), Inches(0.18), _H, theme.primary)
    _rect(slide, Inches(0.18), Emu(0), Inches(0.06), _H, theme.secondary)


def _footer(slide, theme: Theme, deck_title: str, index: int, total: int) -> None:
    _text(
        slide,
        _MARGIN,
        Inches(7.05),
        Inches(9),
        Inches(0.35),
        deck_title,
        font=theme.body_font,
        size=10,
        color=theme.muted,
    )
    _text(
        slide,
        Inches(11.5),
        Inches(7.05),
        Inches(1.4),
        Inches(0.35),
        f"{index} / {total}",
        font=theme.body_font,
        size=10,
        color=theme.muted,
        align=PP_ALIGN.RIGHT,
    )


# --- Per-layout renderers -------------------------------------------------

def _render_title(slide, s: Slide, theme: Theme) -> None:
    # Big colour-block band behind the title.
    _rect(slide, Emu(0), Inches(2.4), _W, Inches(2.7), theme.surface)
    _rect(slide, _MARGIN, Inches(2.55), Inches(1.6), Inches(0.14), theme.primary)
    _text(
        slide, _MARGIN, Inches(2.8), Inches(11.5), Inches(1.7),
        s.title, font=theme.heading_font, size=46, color=theme.text,
        bold=True, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05,
    )
    if s.subtitle:
        _text(
            slide, _MARGIN, Inches(4.5), Inches(11.5), Inches(0.8),
            s.subtitle, font=theme.body_font, size=20, color=theme.primary,
        )
    # Decorative dots bottom-right.
    for i in range(4):
        _rect(
            slide, Inches(11.9 - i * 0.0), Inches(6.1 + 0), Inches(0.0), Inches(0.0),
            theme.secondary, MSO_SHAPE.OVAL,
        )


def _render_section(slide, s: Slide, theme: Theme) -> None:
    _rect(slide, Emu(0), Emu(0), Inches(0.35), _H, theme.primary)
    _text(
        slide, Inches(1.1), Inches(2.7), Inches(11), Inches(1.0),
        s.title, font=theme.heading_font, size=40, color=theme.text, bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if s.subtitle:
        _text(
            slide, Inches(1.1), Inches(3.8), Inches(11), Inches(0.8),
            s.subtitle, font=theme.body_font, size=18, color=theme.muted,
        )


def _render_bullets(slide, s: Slide, theme: Theme) -> None:
    _accent_corner(slide, theme)
    _text(
        slide, _MARGIN, Inches(0.7), Inches(11.5), Inches(1.0),
        s.title, font=theme.heading_font, size=30, color=theme.text, bold=True,
    )
    _rect(slide, _MARGIN, Inches(1.65), Inches(1.2), Inches(0.08), theme.primary)
    if s.subtitle:
        _text(
            slide, _MARGIN, Inches(1.8), Inches(11.5), Inches(0.6),
            s.subtitle, font=theme.body_font, size=16, color=theme.muted,
        )
    _bullets(
        slide, _MARGIN, Inches(2.5), Inches(11.4), Inches(4.2),
        s.bullets, theme, size=20,
    )


def _render_two_column(slide, s: Slide, theme: Theme) -> None:
    _accent_corner(slide, theme)
    _text(
        slide, _MARGIN, Inches(0.7), Inches(11.5), Inches(1.0),
        s.title, font=theme.heading_font, size=30, color=theme.text, bold=True,
    )
    _rect(slide, _MARGIN, Inches(1.65), Inches(1.2), Inches(0.08), theme.primary)

    col_w = Inches(5.5)
    # Left card
    _rect(slide, _MARGIN, Inches(2.2), col_w, Inches(4.3), theme.surface)
    _bullets(
        slide, Inches(1.2), Inches(2.5), Inches(4.9), Inches(3.8),
        s.bullets, theme, size=17,
    )
    # Right card
    right_x = Inches(7.1)
    _rect(slide, right_x, Inches(2.2), col_w, Inches(4.3), theme.surface)
    _bullets(
        slide, Inches(7.4), Inches(2.5), Inches(4.9), Inches(3.8),
        s.column_b or s.bullets, theme, size=17,
    )


def _render_quote(slide, s: Slide, theme: Theme) -> None:
    _rect(slide, Emu(0), Emu(0), _W, _H, theme.surface)
    _rect(slide, _MARGIN, Inches(1.6), Inches(0.9), Inches(0.9), theme.primary)
    _text(
        slide, _MARGIN, Inches(1.4), Inches(2), Inches(1.6),
        "“", font=theme.heading_font, size=96, color=theme.primary, bold=True,
    )
    _text(
        slide, _MARGIN, Inches(2.7), Inches(11.5), Inches(2.6),
        s.quote or s.title, font=theme.heading_font, size=30, color=theme.text,
        bold=True, line_spacing=1.2,
    )
    if s.attribution:
        _text(
            slide, _MARGIN, Inches(5.5), Inches(11.5), Inches(0.7),
            f"— {s.attribution}", font=theme.body_font, size=18,
            color=theme.primary,
        )


def _render_metrics(slide, s: Slide, theme: Theme) -> None:
    _accent_corner(slide, theme)
    _text(
        slide, _MARGIN, Inches(0.7), Inches(11.5), Inches(1.0),
        s.title, font=theme.heading_font, size=30, color=theme.text, bold=True,
    )
    _rect(slide, _MARGIN, Inches(1.65), Inches(1.2), Inches(0.08), theme.primary)

    metrics = s.metrics[:3]
    if not metrics:
        _bullets(slide, _MARGIN, Inches(2.5), Inches(11.4), Inches(4),
                 s.bullets, theme, size=20)
        return

    gap = Inches(0.4)
    total_w = Inches(11.5)
    card_w = Emu(int((total_w - gap * (len(metrics) - 1)) / len(metrics)))
    x = _MARGIN
    accents = [theme.primary, theme.secondary, theme.primary]
    for i, m in enumerate(metrics):
        _rect(slide, x, Inches(2.6), card_w, Inches(3.2), theme.surface)
        _rect(slide, x, Inches(2.6), card_w, Inches(0.14), accents[i % len(accents)])
        _text(
            slide, x, Inches(3.1), card_w, Inches(1.5),
            m.value, font=theme.heading_font, size=48,
            color=accents[i % len(accents)], bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        _text(
            slide, x, Inches(4.7), card_w, Inches(1.0),
            m.label, font=theme.body_font, size=16, color=theme.text,
            align=PP_ALIGN.CENTER,
        )
        x = Emu(x + card_w + gap)


def _render_closing(slide, s: Slide, theme: Theme) -> None:
    _rect(slide, Emu(0), Emu(0), _W, _H, theme.surface)
    _rect(slide, Inches(5.6), Inches(2.5), Inches(2.1), Inches(0.14), theme.primary)
    _text(
        slide, _MARGIN, Inches(2.8), Inches(11.5), Inches(1.4),
        s.title or "Thank You", font=theme.heading_font, size=44,
        color=theme.text, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if s.subtitle:
        _text(
            slide, _MARGIN, Inches(4.3), Inches(11.5), Inches(0.8),
            s.subtitle, font=theme.body_font, size=20, color=theme.primary,
            align=PP_ALIGN.CENTER,
        )


_RENDERERS = {
    SlideLayout.TITLE: _render_title,
    SlideLayout.SECTION: _render_section,
    SlideLayout.BULLETS: _render_bullets,
    SlideLayout.TWO_COLUMN: _render_two_column,
    SlideLayout.QUOTE: _render_quote,
    SlideLayout.METRICS: _render_metrics,
    SlideLayout.CLOSING: _render_closing,
}


def render_deck(deck: Deck, output_dir: str | Path) -> Path:
    """Render ``deck`` to a .pptx file and return its path."""
    theme = get_theme(deck.theme)
    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H

    total = len(deck.slides)
    chromeless = {SlideLayout.TITLE, SlideLayout.SECTION,
                  SlideLayout.QUOTE, SlideLayout.CLOSING}

    for idx, s in enumerate(deck.slides, start=1):
        slide = _blank_slide(prs)
        _fill_background(slide, theme.background)
        _RENDERERS.get(s.layout, _render_bullets)(slide, s, theme)
        if s.layout not in chromeless:
            _footer(slide, theme, deck.title, idx, total)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in " -_" else "" for c in deck.title)
    safe = "-".join(safe.split()).lower()[:60] or "deck"
    path = out_dir / f"{safe}.pptx"
    # Avoid clobbering an existing deck with the same title.
    n = 2
    while path.exists():
        path = out_dir / f"{safe}-{n}.pptx"
        n += 1
    prs.save(str(path))
    return path
